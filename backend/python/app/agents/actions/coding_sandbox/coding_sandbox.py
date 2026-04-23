"""CodingSandbox toolset -- execute Python and TypeScript code in an isolated sandbox.

Exposes ``execute_python`` and ``execute_typescript`` tools to the agent.
Produced artifacts (files written to the output directory) are uploaded to
blob storage and surfaced as ``::artifact`` markers in the response.
"""

from __future__ import annotations

import asyncio
import json
import logging
import os
from pathlib import Path
from typing import Any, Literal, Optional

from pydantic import BaseModel, Field

from app.agents.tools.config import ToolCategory
from app.agents.tools.decorator import tool
from app.agents.tools.models import ToolIntent
from app.connectors.core.registry.auth_builder import AuthBuilder
from app.connectors.core.registry.tool_builder import (
    ToolsetBuilder,
    ToolsetCategory,
)
from app.modules.agents.qna.chat_state import ChatState
from app.sandbox.artifact_upload import upload_artifacts_to_blob
from app.sandbox.manager import get_executor
from app.sandbox.models import ArtifactOutput, ExecutionResult, SandboxLanguage
from app.sandbox.package_policy import PackageNotAllowedError, get_allowlist
from app.sandbox.redact import redact_sandbox_paths
from app.utils.conversation_tasks import register_task

logger = logging.getLogger(__name__)

# Pydantic schemas for tool inputs

class ExecutePythonInput(BaseModel):
    code: str = Field(
        ...,
        description=(
            "Python code to execute. "
            "IMPORTANT: The sandbox has NO internet access -- do not use urllib, requests, httpx, or any network calls. "
            "All data must be passed inline in the code or generated within the code itself. "
            "Write output files to the path in the OUTPUT_DIR environment variable. "
            "Example: open(os.path.join(os.environ['OUTPUT_DIR'], 'chart.png'), 'wb')"
        ),
    )
    requirements: list[str] = Field(
        default_factory=list,
        description=(
            "pip packages to install before execution (e.g. ['pandas', 'matplotlib']). "
            "Only packages on the sandbox allowlist are accepted; any other package will be rejected. "
            "Allowlist includes: pandas, numpy, scipy, scikit-learn, matplotlib, seaborn, "
            "pillow, openpyxl, python-docx, python-pptx, reportlab, pypdf, pdfplumber, plotly, "
            "kaleido, fpdf2, cairosvg, beautifulsoup4, lxml, jinja2, and a curated set more."
        ),
    )


class GetDocumentSkillInput(BaseModel):
    kind: Literal["pptx", "docx"] = Field(
        ...,
        description=(
            "Which design skill to load: 'pptx' for slide decks / presentations, "
            "'docx' for Word documents / reports."
        ),
    )


class IngestTemplateInput(BaseModel):
    file_name: str = Field(
        ...,
        description=(
            "Name of the .pptx or .docx template file to inspect. Must be an "
            "artifact produced by (or uploaded into) a prior sandbox execution "
            "in this conversation. The tool returns a structured summary of the "
            "template's layouts, styles, and theme colors so the subsequent "
            "execute_python call can preserve the template's branding."
        ),
    )


class _SpecSlide(BaseModel):
    """One slide in a ``PresentationSpec``.

    The discriminator is ``type``; the remaining fields are a union of every
    primitive's options. Pydantic validates that the right combination is
    present via ``model_validator`` in the enclosing renderer.
    """
    type: Literal[
        "title",
        "section_divider",
        "two_column",
        "stat_grid",
        "icon_rows",
        "timeline",
        "closing",
    ]
    # Title / section / closing
    title: Optional[str] = None
    subtitle: Optional[str] = None
    eyebrow: Optional[str] = None
    cta: Optional[str] = None
    author: Optional[str] = None
    date: Optional[str] = None
    # Two column
    left: Optional[dict] = None
    right: Optional[dict] = None
    split: Optional[int] = None
    # Stat grid
    stats: Optional[list[dict]] = None
    # Icon rows
    rows: Optional[list[dict]] = None
    # Timeline
    steps: Optional[list[dict]] = None


class PresentationSpec(BaseModel):
    """Deterministic, validated description of a slide deck.

    The LLM produces one of these and the server-side renderer turns it into
    a .pptx via the in-repo ``pipeshub-slides`` library. No LLM-authored code
    runs at render time so the output is deterministic and brand-compliant.
    """
    title: Optional[str] = None
    author: Optional[str] = None
    company: Optional[str] = None
    theme: str = Field(
        default="midnightExecutive",
        description=(
            "Theme name from pipeshub-slides' THEMES registry. Valid values: "
            "midnightExecutive, forestMoss, coralEnergy, warmTerracotta, "
            "oceanGradient, charcoalMinimal, tealTrust, berryCream, sageCalm, "
            "cherryBold."
        ),
    )
    layout: Literal["16x9", "16x10", "4x3", "wide"] = "16x9"
    file_name: str = Field(
        default="presentation.pptx",
        description="Output file name (will be written to OUTPUT_DIR).",
    )
    slides: list[_SpecSlide] = Field(
        default_factory=list,
        description="Ordered list of slides. Must contain at least one slide.",
    )


class RenderPresentationFromSpecInput(BaseModel):
    spec: PresentationSpec = Field(
        ...,
        description=(
            "Structured description of the slide deck (title, theme, slides). "
            "Every slide must set `type` to one of 'title', 'section_divider', "
            "'two_column', 'stat_grid', 'icon_rows', 'timeline', 'closing', "
            "and populate the fields that primitive needs."
        ),
    )


class RenderArtifactPreviewInput(BaseModel):
    file_name: str = Field(
        ...,
        description=(
            "Name of the .pptx, .docx, or .xlsx artifact file produced by a "
            "prior execute_python / execute_typescript call. The tool renders "
            "it to one JPEG per slide / page so the agent can visually inspect "
            "the output for overflow, overlap, low contrast, or alignment "
            "issues before declaring the task done."
        ),
    )
    max_pages: int = Field(
        default=20,
        ge=1,
        le=100,
        description=(
            "Maximum number of pages / slides to render. Rendering is capped "
            "to avoid running away on long documents."
        ),
    )


class ExecuteTypeScriptInput(BaseModel):
    code: str = Field(
        ...,
        description=(
            "TypeScript code to execute. "
            "IMPORTANT: The sandbox has NO internet access -- do not use fetch, axios, http, or any network calls. "
            "All data must be passed inline in the code or generated within the code itself. "
            "Write output files to the path in the OUTPUT_DIR environment variable. "
            "Example: fs.writeFileSync(path.join(process.env.OUTPUT_DIR!, 'report.html'), html)"
        ),
    )
    packages: list[str] = Field(
        default_factory=list,
        description=(
            "npm packages to install before execution (e.g. ['chart.js', 'docx']). "
            "Only packages on the sandbox allowlist are accepted; any other package will be rejected. "
            "Allowlist: fs-extra, sharp, @types/node, csv-stringify, json2csv, chart.js, "
            "docx, pdfkit, jsdom, xlsx, papaparse."
        ),
    )


# -------------------------------------------------------------------
# Toolset registration
# -------------------------------------------------------------------

@ToolsetBuilder("Coding Sandbox")\
    .in_group("Internal Tools")\
    .with_description("Execute Python and TypeScript code in a sandboxed environment to generate artifacts like charts, documents, images, and data files")\
    .with_category(ToolsetCategory.CODE_EXECUTION)\
    .with_auth([
        AuthBuilder.type("NONE").fields([])
    ])\
    .as_internal()\
    .configure(lambda builder: builder.with_icon("/assets/icons/toolsets/coding_sandbox.svg"))\
    .build_decorator()
class CodingSandbox:
    """Code execution sandbox exposed to agents."""

    def __init__(self, state: ChatState) -> None:
        self.chat_state = state

    def _result(self, success: bool, payload: dict[str, Any]) -> tuple[bool, str]:
        return success, json.dumps(payload, default=str)

    def _audit_log(
        self,
        tool_name: str,
        language: SandboxLanguage,
        packages: list[str] | None,
        code: str,
    ) -> None:
        """Structured audit line for every code-execution tool invocation."""
        logger.info(
            "[%s] AUDIT | org_id=%s user_id=%s conversation_id=%s language=%s "
            "packages=%s code_length=%d",
            tool_name,
            self.chat_state.get("org_id"),
            self.chat_state.get("user_id"),
            self.chat_state.get("conversation_id"),
            language.value,
            packages or [],
            len(code),
        )
        logger.debug("[%s] CODE:\n%s", tool_name, code)

    def _package_rejection(self, exc: PackageNotAllowedError) -> tuple[bool, str]:
        """Convert a PackageNotAllowedError into a structured tool error.

        The LLM can read ``error_code`` and ``allowed_packages`` to retry the
        call with a valid package name.
        """
        logger.warning(
            "Sandbox package rejected: org_id=%s user_id=%s conversation_id=%s "
            "language=%s package=%s",
            self.chat_state.get("org_id"),
            self.chat_state.get("user_id"),
            self.chat_state.get("conversation_id"),
            exc.language.value,
            exc.package,
        )
        return self._result(False, {
            "success": False,
            "error_code": "package_not_allowed",
            "message": (
                f"Package {exc.package!r} is not on the {exc.language.value} sandbox "
                f"allowlist and will not be installed."
            ),
            "rejected_package": exc.package,
            "language": exc.language.value,
            "allowed_packages": list(get_allowlist(exc.language)),
        })

    # ------------------------------------------------------------------
    # Artifact upload helper
    # ------------------------------------------------------------------

    async def _upload_artifacts(
        self, exec_result: ExecutionResult, *, source_tool: str = "coding_sandbox.execute_python",
    ) -> list[dict[str, Any]]:
        """Upload produced artifacts to blob storage and create ArtifactRecords.

        Thin wrapper around :func:`app.sandbox.artifact_upload.upload_artifacts_to_blob`
        that resolves context (conversation / org / user / stores) from
        ``chat_state`` and returns the list of upload-info dicts.
        """
        if not exec_result.artifacts:
            return []

        conversation_id = self.chat_state.get("conversation_id")
        org_id = self.chat_state.get("org_id")
        user_id = self.chat_state.get("user_id")
        blob_store = self.chat_state.get("blob_store")
        graph_provider = self.chat_state.get("graph_provider")

        if blob_store is None:
            try:
                from app.modules.transformers.blob_storage import BlobStorage
                blob_store = BlobStorage(
                    logger=logger,
                    config_service=self.chat_state.get("config_service"),
                    graph_provider=graph_provider,
                )
            except (ImportError, OSError, RuntimeError, ValueError) as e:
                logger.warning("Could not init BlobStorage for artifact upload: %s", e)
                return []

        if not (conversation_id and org_id and blob_store):
            logger.warning(
                "Artifact upload skipped: conversation_id=%s org_id=%s blob_store=%s",
                conversation_id, org_id, "yes" if blob_store else "no",
            )
            return []

        from app.config.constants.arangodb import Connectors

        return await upload_artifacts_to_blob(
            exec_result.artifacts,
            blob_store=blob_store,
            org_id=org_id,
            conversation_id=conversation_id,
            user_id=user_id,
            graph_provider=graph_provider,
            connector_name=Connectors.CODING_SANDBOX,
            source_tool=source_tool,
        )

    def _schedule_artifact_upload(
        self, exec_result: ExecutionResult, *, source_tool: str = "coding_sandbox.execute_python",
    ) -> None:
        """Schedule artifact upload as a background conversation task."""
        conversation_id = self.chat_state.get("conversation_id")
        if not conversation_id or not exec_result.artifacts:
            return

        async def _upload() -> Optional[dict[str, Any]]:
            try:
                results = await self._upload_artifacts(exec_result, source_tool=source_tool)
                if results:
                    return {"type": "artifacts", "artifacts": results}
            except (OSError, ConnectionError, RuntimeError, ValueError):
                logger.exception("Background artifact upload failed")
            return None

        task = asyncio.create_task(_upload())
        register_task(conversation_id, task)

    # ------------------------------------------------------------------
    # Tools
    # ------------------------------------------------------------------

    @tool(
        app_name="coding_sandbox",
        tool_name="execute_python",
        args_schema=ExecutePythonInput,
        llm_description=(
            "Execute Python code in an isolated sandbox with NO internet access. "
            "The sandbox cannot make any network requests (no urllib, requests, httpx, API calls, web scraping, etc.). "
            "All data needed for computation must be embedded directly in the code as literals or variables. "
            "Use this to generate files such as charts (matplotlib/plotly), documents (python-docx/python-pptx), "
            "images (Pillow), CSV/Excel (pandas/openpyxl), PDF (reportlab/fpdf), or perform data processing. "
            "Write output files to os.path.join(os.environ['OUTPUT_DIR'], 'filename'). "
            "Common libraries are pre-installed: pandas, matplotlib, seaborn, plotly, openpyxl, "
            "python-docx, python-pptx, Pillow, reportlab, fpdf2, cairosvg. "
            "RICH DOCUMENTS: For .pptx/.docx output that should be visually polished (not a wall of "
            "bullets on a white background), FIRST call coding_sandbox.get_document_skill with kind='pptx' "
            "or kind='docx' to load the curated design system (palettes, typography, layouts, pitfalls). "
            "For new decks/docs from scratch, prefer execute_typescript with the `pipeshub-slides` / "
            "`pipeshub-docs` helper libraries (far better styling than python-pptx/python-docx). "
            "Use execute_python with python-pptx/python-docx when editing an uploaded template, since "
            "that preserves masters and themes. After writing a .pptx/.docx, call "
            "coding_sandbox.render_artifact_preview on it to visually inspect the output before declaring done. "
            "IMPORTANT: Do NOT print or echo the output file's absolute path (e.g. do not `print(out_path)` "
            "or `print(f'Saved at: {out_path}')`). Produced files are attached to the response automatically "
            "as artifacts; the UI renders download cards for them. Any file path printed to stdout is an "
            "internal sandbox path and must NEVER appear in your answer to the user."
        ),
        category=ToolCategory.CODE_EXECUTION,
        is_essential=False,
        requires_auth=False,
        when_to_use=[
            "User wants to generate a chart, graph, or visualization",
            "User wants to create a document (DOCX, PPTX, PDF)",
            "User wants to process data and produce output files",
            "User wants to generate an image or transform images",
            "User wants to create a CSV or Excel file from data",
            "User needs to run a computation that produces file output",
        ],
        when_not_to_use=[
            "User just wants a text answer or explanation",
            "User wants to query an existing database (use database tools instead)",
            "User wants to search knowledge base (use retrieval tools)",
            "The task requires downloading data from the internet or calling external APIs (the sandbox has no network access)",
        ],
        primary_intent=ToolIntent.ACTION,
        typical_queries=[
            "Create a bar chart of quarterly revenue",
            "Generate a DOCX report summarizing the data",
            "Make a pie chart showing market share",
            "Create an Excel spreadsheet with the results",
            "Generate a PDF invoice",
        ],
    )
    async def execute_python(self, code: str, requirements: list[str] | None = None) -> tuple[bool, str]:
        """Execute Python code in the sandbox and return results + artifacts."""
        self._audit_log("execute_python", SandboxLanguage.PYTHON, requirements, code)
        # Validate packages at the tool boundary so we can surface a rich,
        # retry-friendly error to the LLM. Executors ALSO re-validate as
        # defence in depth.
        try:
            from app.sandbox.models import validate_packages
            validate_packages(requirements, language=SandboxLanguage.PYTHON)
        except PackageNotAllowedError as exc:
            return self._package_rejection(exc)
        except ValueError as exc:
            return self._result(False, {
                "success": False,
                "error_code": "invalid_package_name",
                "message": str(exc),
            })

        try:
            executor = get_executor()
            logger.debug("[execute_python] executor=%s", type(executor).__name__)
            result = await executor.execute(
                code=code,
                language=SandboxLanguage.PYTHON,
                packages=requirements or [],
            )
            logger.info(
                "[execute_python] OUTPUT | success=%s exit_code=%s time_ms=%s artifacts=%d "
                "stdout_len=%d stderr_len=%d error=%s",
                result.success, result.exit_code, result.execution_time_ms,
                len(result.artifacts), len(result.stdout), len(result.stderr),
                result.error or "None",
            )
            if result.stdout:
                logger.debug("[execute_python] STDOUT:\n%s", _truncate(result.stdout, 2000))
            if result.stderr:
                log_fn = logger.warning if not result.success else logger.debug
                log_fn("[execute_python] STDERR:\n%s", _truncate(result.stderr, 2000))

            self._schedule_artifact_upload(result)

            artifact_info = [
                {"fileName": a.file_name, "mimeType": a.mime_type, "sizeBytes": a.size_bytes}
                for a in result.artifacts
            ]

            return self._result(result.success, {
                "success": result.success,
                "message": "Code executed successfully" if result.success else "Code execution failed",
                "stdout": redact_sandbox_paths(_truncate(result.stdout, 8000)),
                "stderr": redact_sandbox_paths(_truncate(result.stderr, 4000)),
                "exit_code": result.exit_code,
                "execution_time_ms": result.execution_time_ms,
                "artifacts": artifact_info,
                **({"error": redact_sandbox_paths(result.error)} if result.error else {}),
            })
        except Exception as e:
            logger.exception("[execute_python] EXCEPTION")
            return self._result(False, {"success": False, "error": redact_sandbox_paths(str(e))})

    @tool(
        app_name="coding_sandbox",
        tool_name="execute_typescript",
        args_schema=ExecuteTypeScriptInput,
        llm_description=(
            "Execute TypeScript code in an isolated sandbox with NO internet access. "
            "The sandbox cannot make any network requests (no fetch, axios, http, API calls, web scraping, etc.). "
            "All data needed for computation must be embedded directly in the code as literals or variables. "
            "Use this when JavaScript/TypeScript is more appropriate for the task. "
            "Write output files to path.join(process.env.OUTPUT_DIR!, 'filename'). "
            "Common packages are pre-installed: fs-extra, sharp, chart.js, pptxgenjs, docx, "
            "react-icons, and the in-repo helper libraries `pipeshub-slides` and `pipeshub-docs`. "
            "RICH DOCUMENTS: This is the RECOMMENDED path for generating new .pptx / .docx from "
            "scratch. Before writing code, call coding_sandbox.get_document_skill with kind='pptx' "
            "(for slide decks) or kind='docx' (for Word docs) to load the design system (palettes, "
            "typography, layouts, pitfalls, example snippets). Prefer `pipeshub-slides` / "
            "`pipeshub-docs` over raw pptxgenjs/docx-js: the helpers bake in theme, margins, and "
            "styling so you can't produce low-contrast or misaligned output by accident. After "
            "writing a .pptx/.docx, call coding_sandbox.render_artifact_preview on it to visually "
            "inspect the output and iterate on any issues before declaring done. "
            "IMPORTANT: Do NOT print or echo the output file's absolute path (e.g. do not "
            "`console.log(outPath)` or `console.log(`Saved at: ${outPath}`)`). Produced files are "
            "attached to the response automatically as artifacts; the UI renders download cards for "
            "them. Any file path printed to stdout is an internal sandbox path and must NEVER appear "
            "in your answer to the user."
        ),
        category=ToolCategory.CODE_EXECUTION,
        is_essential=False,
        requires_auth=False,
        when_to_use=[
            "User needs TypeScript/JavaScript-specific processing",
            "User wants to generate HTML/SVG visualizations",
            "User wants to use Node.js libraries for file generation",
        ],
        when_not_to_use=[
            "User just wants a text answer",
            "Python would be more suitable for the task",
            "The task requires downloading data from the internet or calling external APIs (the sandbox has no network access)",
        ],
        primary_intent=ToolIntent.ACTION,
        typical_queries=[
            "Generate an HTML report",
            "Create an SVG diagram",
        ],
    )
    async def execute_typescript(self, code: str, packages: list[str] | None = None) -> tuple[bool, str]:
        """Execute TypeScript code in the sandbox and return results + artifacts."""
        self._audit_log("execute_typescript", SandboxLanguage.TYPESCRIPT, packages, code)
        try:
            from app.sandbox.models import validate_packages
            validate_packages(packages, language=SandboxLanguage.TYPESCRIPT)
        except PackageNotAllowedError as exc:
            return self._package_rejection(exc)
        except ValueError as exc:
            return self._result(False, {
                "success": False,
                "error_code": "invalid_package_name",
                "message": str(exc),
            })

        try:
            executor = get_executor()
            logger.debug("[execute_typescript] executor=%s", type(executor).__name__)
            result = await executor.execute(
                code=code,
                language=SandboxLanguage.TYPESCRIPT,
                packages=packages or [],
            )
            logger.info(
                "[execute_typescript] OUTPUT | success=%s exit_code=%s time_ms=%s artifacts=%d "
                "stdout_len=%d stderr_len=%d error=%s",
                result.success, result.exit_code, result.execution_time_ms,
                len(result.artifacts), len(result.stdout), len(result.stderr),
                result.error or "None",
            )
            if result.stdout:
                logger.debug("[execute_typescript] STDOUT:\n%s", _truncate(result.stdout, 2000))
            if result.stderr:
                log_fn = logger.warning if not result.success else logger.debug
                log_fn("[execute_typescript] STDERR:\n%s", _truncate(result.stderr, 2000))

            self._schedule_artifact_upload(result, source_tool="coding_sandbox.execute_typescript")

            artifact_info = [
                {"fileName": a.file_name, "mimeType": a.mime_type, "sizeBytes": a.size_bytes}
                for a in result.artifacts
            ]

            return self._result(result.success, {
                "success": result.success,
                "message": "Code executed successfully" if result.success else "Code execution failed",
                "stdout": redact_sandbox_paths(_truncate(result.stdout, 8000)),
                "stderr": redact_sandbox_paths(_truncate(result.stderr, 4000)),
                "exit_code": result.exit_code,
                "execution_time_ms": result.execution_time_ms,
                "artifacts": artifact_info,
                **({"error": redact_sandbox_paths(result.error)} if result.error else {}),
            })
        except Exception as e:
            logger.exception("[execute_typescript] EXCEPTION")
            return self._result(False, {"success": False, "error": redact_sandbox_paths(str(e))})

    # ------------------------------------------------------------------
    # Document design skill
    # ------------------------------------------------------------------

    @tool(
        app_name="coding_sandbox",
        tool_name="get_document_skill",
        args_schema=GetDocumentSkillInput,
        llm_description=(
            "Load the document design skill — a curated guide for producing visually polished "
            ".pptx / .docx artifacts (color palettes, typography, layout catalogue, pitfalls, "
            "example snippets). Call this BEFORE writing any code that generates a slide deck or "
            "Word document so the produced artifact has structure and styling instead of being a "
            "wall of bullets on a white background. Returns the skill content as a string in the "
            "'skill' field; read it carefully and follow its guidance in the subsequent "
            "execute_python / execute_typescript call."
        ),
        category=ToolCategory.CODE_EXECUTION,
        is_essential=False,
        requires_auth=False,
        when_to_use=[
            "User asks to create, generate, or produce a .pptx / slide deck / pitch deck / "
            "presentation",
            "User asks to create, generate, or produce a .docx / Word document / report / memo / "
            "brief / letter / proposal",
            "User asks to edit or customize an existing .pptx or .docx",
        ],
        when_not_to_use=[
            "User wants a chart or image only (no slide deck)",
            "User wants a PDF, CSV, or XLSX (use the relevant guidance instead)",
            "Pure text answer",
        ],
        primary_intent=ToolIntent.QUESTION,
        typical_queries=[
            "How should I style this slide deck?",
            "Load the document design guide before I make the report",
        ],
    )
    async def get_document_skill(self, kind: str) -> tuple[bool, str]:
        """Return the curated design skill (markdown) for .pptx or .docx output."""
        if kind not in {"pptx", "docx"}:
            return self._result(False, {
                "success": False,
                "error_code": "invalid_kind",
                "message": f"kind must be 'pptx' or 'docx', got {kind!r}",
            })

        try:
            skill = _load_skill_markdown(kind)
        except FileNotFoundError as exc:
            logger.exception("[get_document_skill] could not locate %s skill", kind)
            return self._result(False, {
                "success": False,
                "error_code": "skill_missing",
                "message": f"Could not load the {kind} design skill: {exc}",
            })

        logger.info(
            "[get_document_skill] AUDIT | org_id=%s user_id=%s conversation_id=%s kind=%s bytes=%d",
            self.chat_state.get("org_id"),
            self.chat_state.get("user_id"),
            self.chat_state.get("conversation_id"),
            kind,
            len(skill),
        )
        return self._result(True, {
            "success": True,
            "kind": kind,
            "skill": skill,
            "message": (
                f"Loaded the {kind} design skill. Follow its guidance (palettes, typography, "
                f"layouts, pitfalls) in your subsequent execute_python / execute_typescript call."
            ),
        })

    # ------------------------------------------------------------------
    # Template ingestion
    # ------------------------------------------------------------------

    @tool(
        app_name="coding_sandbox",
        tool_name="ingest_template",
        args_schema=IngestTemplateInput,
        llm_description=(
            "Inspect an existing .pptx or .docx template file (from a prior sandbox execution in "
            "this conversation) and return a structured summary of its slide layouts / paragraph "
            "styles / theme colors. Call this BEFORE writing code that edits the template so the "
            "edits preserve the user's branding. The tool also returns `template_path` — the "
            "absolute path on the sandbox host where the template lives — which you can then pass "
            "into a subsequent execute_python call via Presentation(template_path) or "
            "Document(template_path)."
        ),
        category=ToolCategory.CODE_EXECUTION,
        is_essential=False,
        requires_auth=False,
        when_to_use=[
            "User uploaded a .pptx/.docx template and asked to fill it in / customize it",
            "You need to preserve an existing deck's master slides / theme colors before editing",
        ],
        when_not_to_use=[
            "Creating a document from scratch (no template)",
            "The file is not a .pptx or .docx",
        ],
        primary_intent=ToolIntent.ANALYSIS,
        typical_queries=[
            "Use this template to build the report",
            "Fill this deck with these bullets",
        ],
    )
    async def ingest_template(self, file_name: str) -> tuple[bool, str]:
        """Locate a .pptx/.docx file in sandbox dirs and return a structural summary."""
        logger.info(
            "[ingest_template] AUDIT | org_id=%s user_id=%s conversation_id=%s file=%s",
            self.chat_state.get("org_id"),
            self.chat_state.get("user_id"),
            self.chat_state.get("conversation_id"),
            file_name,
        )

        match = _find_latest_sandbox_file(file_name)
        if match is None:
            return self._result(False, {
                "success": False,
                "error_code": "file_not_found",
                "message": (
                    f"Could not locate {file_name!r} in any recent sandbox directory. The template "
                    f"must be an artifact from a prior sandbox execution in this conversation."
                ),
            })

        ext = match.suffix.lower()
        if ext == ".pptx":
            summary = _summarise_pptx(match)
        elif ext == ".docx":
            summary = _summarise_docx(match)
        else:
            return self._result(False, {
                "success": False,
                "error_code": "unsupported_extension",
                "message": f"ingest_template only supports .pptx and .docx (got {ext!r}).",
            })

        return self._result(True, {
            "success": True,
            "file_name": file_name,
            "template_path": str(match),
            "summary": summary,
            "message": (
                f"Inspected {file_name}. Pass `template_path` into your next execute_python call "
                f"(e.g. Presentation(template_path) or Document(template_path)) to edit the "
                f"template while preserving its masters and styles."
            ),
        })

    # ------------------------------------------------------------------
    # Deterministic JSON -> pptx renderer (v2)
    # ------------------------------------------------------------------

    @tool(
        app_name="coding_sandbox",
        tool_name="render_presentation_from_spec",
        args_schema=RenderPresentationFromSpecInput,
        llm_description=(
            "Deterministic alternative to execute_typescript for producing a .pptx. Provide a "
            "structured PresentationSpec (title, theme, slides list) and the server renders it "
            "into a slide deck using the pipeshub-slides helper — no LLM-authored code runs at "
            "render time, so the output is fully deterministic and brand-compliant. Use this "
            "when you have clear bullet points / stats / sections in mind and don't need custom "
            "layout code. Fall back to execute_typescript if you need anything the spec doesn't "
            "cover (custom charts, images, exotic layouts). Each slide's `type` must be one of: "
            "'title', 'section_divider', 'two_column', 'stat_grid', 'icon_rows', 'timeline', "
            "'closing'. See `get_document_skill(kind='pptx')` for the design system this renders "
            "against."
        ),
        category=ToolCategory.CODE_EXECUTION,
        is_essential=False,
        requires_auth=False,
        when_to_use=[
            "User supplies clear, structured content (bullets, stats, sections) for a deck",
            "Previous execute_typescript attempts produced broken layouts",
            "You want deterministic, brand-compliant output without running LLM code",
        ],
        when_not_to_use=[
            "You need custom images, charts with unusual shapes, or one-off layouts",
            "The deck must be produced from an uploaded template (use ingest_template + execute_python)",
        ],
        primary_intent=ToolIntent.ACTION,
        typical_queries=[
            "Make me a deck from these bullets",
            "Build a quick status update slide deck",
        ],
    )
    async def render_presentation_from_spec(
        self, spec: dict[str, Any] | PresentationSpec,
    ) -> tuple[bool, str]:
        """Render a PresentationSpec JSON into a .pptx deterministically."""
        try:
            validated = (
                spec if isinstance(spec, PresentationSpec)
                else PresentationSpec.model_validate(spec)
            )
        except Exception as exc:
            return self._result(False, {
                "success": False,
                "error_code": "invalid_spec",
                "message": f"Invalid PresentationSpec: {exc}",
            })

        if not validated.slides:
            return self._result(False, {
                "success": False,
                "error_code": "empty_spec",
                "message": "PresentationSpec must contain at least one slide.",
            })

        spec_json = validated.model_dump_json()
        safe_file_name = os.path.basename(validated.file_name or "presentation.pptx")
        if not safe_file_name.lower().endswith(".pptx"):
            safe_file_name = safe_file_name + ".pptx"

        code = _render_from_spec_code(spec_json, safe_file_name)

        logger.info(
            "[render_presentation_from_spec] AUDIT | org_id=%s user_id=%s "
            "conversation_id=%s slides=%d theme=%s",
            self.chat_state.get("org_id"),
            self.chat_state.get("user_id"),
            self.chat_state.get("conversation_id"),
            len(validated.slides),
            validated.theme,
        )

        try:
            executor = get_executor()
            result = await executor.execute(
                code=code,
                language=SandboxLanguage.TYPESCRIPT,
                packages=["pptxgenjs"],
            )
        except PackageNotAllowedError as exc:
            return self._result(False, {
                "success": False,
                "error_code": "package_not_allowed",
                "message": str(exc),
            })
        except Exception as exc:
            logger.exception("[render_presentation_from_spec] executor failed")
            return self._result(False, {
                "success": False,
                "error_code": "executor_failed",
                "message": redact_sandbox_paths(str(exc)),
            })

        self._schedule_artifact_upload(
            result, source_tool="coding_sandbox.render_presentation_from_spec",
        )

        artifact_info = [
            {"fileName": a.file_name, "mimeType": a.mime_type, "sizeBytes": a.size_bytes}
            for a in result.artifacts
        ]
        return self._result(result.success, {
            "success": result.success,
            "file_name": safe_file_name,
            "slide_count": len(validated.slides),
            "theme": validated.theme,
            "artifacts": artifact_info,
            "stderr": redact_sandbox_paths(_truncate(result.stderr or "", 2000)),
            "message": (
                f"Rendered {len(validated.slides)} slides into {safe_file_name} using the "
                f"'{validated.theme}' theme."
                if result.success
                else "Rendering failed — inspect stderr."
            ),
            **({"error": redact_sandbox_paths(result.error)} if result.error else {}),
        })

    # ------------------------------------------------------------------
    # Visual QA -- render pptx/docx to images
    # ------------------------------------------------------------------

    @tool(
        app_name="coding_sandbox",
        tool_name="render_artifact_preview",
        args_schema=RenderArtifactPreviewInput,
        llm_description=(
            "Render a .pptx / .docx / .xlsx artifact (produced by a prior execute_python / "
            "execute_typescript call) to one JPEG per slide or page, so the agent can visually "
            "inspect the output. Use this as the final QA step before declaring a document task "
            "done — the first render is almost never correct. After viewing the images, look for "
            "text overflow, overlapping elements, low-contrast text, misaligned columns, wrapped "
            "titles colliding with body content, and leftover placeholder text. If any issue is "
            "found, fix it in a follow-up execute_python / execute_typescript call and re-render. "
            "Cap at 3 iterations per task. Requires LibreOffice + poppler to be installed on the "
            "sandbox host; disabled if SANDBOX_VISUAL_QA=off."
        ),
        category=ToolCategory.CODE_EXECUTION,
        is_essential=False,
        requires_auth=False,
        when_to_use=[
            "After producing a .pptx / .docx artifact to verify it is visually correct",
            "Before telling the user a deck / report is done",
        ],
        when_not_to_use=[
            "The artifact is not a .pptx / .docx / .xlsx",
            "The agent already rendered and inspected this exact artifact once",
            "Operator has disabled visual QA via SANDBOX_VISUAL_QA=off",
        ],
        primary_intent=ToolIntent.ANALYSIS,
        typical_queries=[
            "Check the deck I just generated",
            "Show me what the report looks like",
        ],
    )
    async def render_artifact_preview(
        self, file_name: str, max_pages: int = 20,
    ) -> tuple[bool, str]:
        """Render a document artifact to JPEG page images and upload them."""
        if os.environ.get("SANDBOX_VISUAL_QA", "on").lower() in {"off", "false", "0", "no"}:
            return self._result(False, {
                "success": False,
                "error_code": "visual_qa_disabled",
                "message": (
                    "Visual QA is disabled on this sandbox "
                    "(SANDBOX_VISUAL_QA=off). Inspect the artifact manually."
                ),
            })

        logger.info(
            "[render_artifact_preview] AUDIT | org_id=%s user_id=%s conversation_id=%s file=%s",
            self.chat_state.get("org_id"),
            self.chat_state.get("user_id"),
            self.chat_state.get("conversation_id"),
            file_name,
        )

        src = _find_latest_sandbox_file(file_name)
        if src is None:
            return self._result(False, {
                "success": False,
                "error_code": "file_not_found",
                "message": (
                    f"Could not locate {file_name!r} in any recent sandbox directory. The file "
                    f"must be an artifact from a prior sandbox execution in this conversation."
                ),
            })

        ext = src.suffix.lower()
        if ext not in {".pptx", ".docx", ".xlsx", ".pdf"}:
            return self._result(False, {
                "success": False,
                "error_code": "unsupported_extension",
                "message": (
                    f"render_artifact_preview only supports .pptx, .docx, .xlsx, .pdf "
                    f"(got {ext!r})."
                ),
            })

        try:
            images, warnings = await asyncio.to_thread(
                _render_to_images, src, max_pages,
            )
        except FileNotFoundError as exc:
            return self._result(False, {
                "success": False,
                "error_code": "renderer_missing",
                "message": (
                    f"Renderer not available on this host: {exc}. Install LibreOffice "
                    f"(`soffice`) and poppler (`pdftoppm`) or disable visual QA via "
                    f"SANDBOX_VISUAL_QA=off."
                ),
            })
        except Exception as exc:
            logger.exception("[render_artifact_preview] render failed for %s", src)
            return self._result(False, {
                "success": False,
                "error_code": "render_failed",
                "message": f"Failed to render {file_name}: {exc}",
            })

        # Build synthetic ArtifactOutput objects and reuse the existing
        # upload + streaming plumbing so the page images surface in the UI
        # the same way any other artifact does.
        artifacts = [
            ArtifactOutput(
                file_name=img.name,
                file_path=str(img),
                mime_type="image/jpeg",
                size_bytes=img.stat().st_size if img.exists() else 0,
            )
            for img in images
        ]
        exec_result = ExecutionResult(success=True, artifacts=artifacts)
        self._schedule_artifact_upload(
            exec_result, source_tool="coding_sandbox.render_artifact_preview",
        )

        artifact_info = [
            {"fileName": a.file_name, "mimeType": a.mime_type, "sizeBytes": a.size_bytes}
            for a in artifacts
        ]
        return self._result(True, {
            "success": True,
            "file_name": file_name,
            "page_count": len(images),
            "artifacts": artifact_info,
            "warnings": warnings,
            "message": (
                f"Rendered {len(images)} page(s) of {file_name}. Inspect the attached images for "
                f"overflow, overlap, low contrast, misalignment, or leftover placeholder text. "
                f"If any issue is found, fix it in a follow-up sandbox call and re-render."
            ),
        })


# ------------------------------------------------------------------
# Helpers
# ------------------------------------------------------------------

def _truncate(text: str, max_len: int) -> str:
    if len(text) <= max_len:
        return text
    return text[:max_len] + f"\n... (truncated, {len(text)} total chars)"


def _load_skill_markdown(kind: str) -> str:
    """Return the raw markdown for the ``pptx`` or ``docx`` design skill.

    We try every known packaging layout so the same code path works for:

    * local source checkouts — ``Path(__file__).parent / "skills"``
    * pip-installed wheels — ``importlib.resources.files(<pkg>) / "skills"``
    * PyInstaller frozen binaries — ``sys._MEIPASS`` extraction dir

    Raises ``FileNotFoundError`` with a joined list of every path we probed
    if the file can't be located anywhere. That keeps the error message
    actionable ("none of these locations contained pptx.md") rather than
    pointing at a single stale path that happens to not exist on this host.
    """
    import sys

    filename = f"{kind}.md"
    attempted: list[str] = []

    # 1. Source-checkout path: the file is a sibling of this module.
    direct = Path(__file__).parent / "skills" / filename
    attempted.append(str(direct))
    if direct.is_file():
        return direct.read_text(encoding="utf-8")

    # 2. importlib.resources — works both for normal installs and for
    #    editable installs when package-data is declared correctly.
    try:
        from importlib.resources import files as _pkg_files
        resource = _pkg_files(
            "app.agents.actions.coding_sandbox",
        ).joinpath("skills", filename)
        attempted.append(str(resource))
        # ``is_file`` exists on Traversable / MultiplexedPath; fall back
        # to read on failure.
        try:
            if resource.is_file():  # type: ignore[attr-defined]
                return resource.read_text(encoding="utf-8")
        except (AttributeError, OSError):
            pass
    except (ImportError, ModuleNotFoundError, TypeError):
        pass

    # 3. PyInstaller: bundled data is unpacked under sys._MEIPASS.
    meipass = getattr(sys, "_MEIPASS", None)
    if meipass:
        frozen = (
            Path(meipass) / "app" / "agents" / "actions"
            / "coding_sandbox" / "skills" / filename
        )
        attempted.append(str(frozen))
        if frozen.is_file():
            return frozen.read_text(encoding="utf-8")

    raise FileNotFoundError(
        f"{filename} not found; probed: " + "; ".join(attempted),
    )


# Roots inspected by render_artifact_preview / ingest_template. Kept as a
# module-level tuple so tests can monkey-patch if needed.
def _sandbox_roots() -> tuple[Path, ...]:
    # Import lazily to avoid circular dependencies at module-import time.
    from app.sandbox.docker_executor import DockerExecutor
    from app.sandbox.local_executor import LocalExecutor
    return (Path(LocalExecutor.get_sandbox_root()), Path(DockerExecutor.get_sandbox_root()))


def _find_latest_sandbox_file(file_name: str) -> Path | None:
    """Return the most recently modified match for *file_name* across sandbox roots.

    File name is taken as a literal basename — directory traversal (``..``,
    leading slash) is rejected to avoid escaping the sandbox tree.
    """
    safe_name = os.path.basename(file_name).strip()
    if not safe_name or safe_name != file_name.strip().lstrip("/").split("/")[-1]:
        # Reject inputs that look like paths rather than plain file names.
        return None
    candidates: list[Path] = []
    for root in _sandbox_roots():
        if not root.is_dir():
            continue
        try:
            candidates.extend(root.rglob(safe_name))
        except OSError:
            continue
    if not candidates:
        return None
    candidates.sort(key=lambda p: p.stat().st_mtime if p.exists() else 0, reverse=True)
    return candidates[0]


def _summarise_pptx(path: Path) -> dict[str, Any]:
    """Return a structural summary of a .pptx file (layouts, placeholders, theme)."""
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return {"error": "python-pptx is not installed on the host"}

    try:
        prs = Presentation(str(path))
    except Exception as exc:
        return {"error": f"Could not open pptx: {exc}"}

    layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for ph in layout.placeholders:
            placeholders.append({
                "idx": int(ph.placeholder_format.idx),
                "name": ph.name,
                "type": str(ph.placeholder_format.type),
            })
        layouts.append({
            "index": i,
            "name": layout.name,
            "placeholders": placeholders,
        })

    return {
        "slide_width_emu": prs.slide_width,
        "slide_height_emu": prs.slide_height,
        "slide_count": len(prs.slides),
        "layouts": layouts,
    }


def _summarise_docx(path: Path) -> dict[str, Any]:
    """Return a structural summary of a .docx file (styles, sections)."""
    try:
        from docx import Document  # type: ignore
    except ImportError:
        return {"error": "python-docx is not installed on the host"}

    try:
        doc = Document(str(path))
    except Exception as exc:
        return {"error": f"Could not open docx: {exc}"}

    styles = []
    seen = set()
    for s in doc.styles:
        name = getattr(s, "name", None)
        if not name or name in seen:
            continue
        seen.add(name)
        styles.append({"name": name, "type": str(s.type)})

    sections = []
    for sec in doc.sections:
        sections.append({
            "page_width_emu": sec.page_width,
            "page_height_emu": sec.page_height,
            "left_margin_emu": sec.left_margin,
            "right_margin_emu": sec.right_margin,
            "top_margin_emu": sec.top_margin,
            "bottom_margin_emu": sec.bottom_margin,
        })

    return {
        "paragraph_count": len(doc.paragraphs),
        "style_count": len(styles),
        "styles": styles[:60],  # cap to avoid bloat
        "sections": sections,
    }


def _render_to_images(src: Path, max_pages: int) -> tuple[list[Path], list[str]]:
    """Convert *src* to a PDF via LibreOffice, then split into per-page JPEGs.

    Returns ``(list_of_image_paths, warnings)``. Images are written next to
    the source file in a ``preview_<stem>`` sibling directory so
    ``_schedule_artifact_upload`` can collect them.
    """
    import subprocess

    if shutil_which("soffice") is None:
        raise FileNotFoundError("LibreOffice 'soffice' is not on PATH")
    if shutil_which("pdftoppm") is None:
        raise FileNotFoundError("poppler 'pdftoppm' is not on PATH")

    preview_dir = src.parent / f"preview_{src.stem}"
    preview_dir.mkdir(parents=True, exist_ok=True)

    warnings: list[str] = []
    ext = src.suffix.lower()
    if ext == ".pdf":
        pdf_path = src
    else:
        result = subprocess.run(
            [
                "soffice", "--headless",
                "--convert-to", "pdf",
                "--outdir", str(preview_dir),
                str(src),
            ],
            capture_output=True, text=True, timeout=180,
        )
        if result.returncode != 0:
            raise RuntimeError(
                f"soffice failed (exit {result.returncode}): {result.stderr.strip()[:500]}"
            )
        pdf_path = preview_dir / (src.stem + ".pdf")
        if not pdf_path.exists():
            raise RuntimeError(f"LibreOffice produced no PDF at {pdf_path}")

    out_prefix = preview_dir / f"{src.stem}-page"
    subprocess.run(
        [
            "pdftoppm", "-jpeg", "-r", "120",
            "-f", "1", "-l", str(max_pages),
            str(pdf_path), str(out_prefix),
        ],
        check=True, timeout=180,
    )

    # pdftoppm pads the page index based on the total page count. Sort
    # naturally so page 2 comes before page 10.
    images = sorted(preview_dir.glob(f"{src.stem}-page-*.jpg"))
    if not images:
        raise RuntimeError("pdftoppm produced no JPEG pages")

    # Clean up intermediate PDF (keep it only when the input itself was a PDF)
    if ext != ".pdf":
        try:
            pdf_path.unlink()
        except OSError:
            warnings.append(f"Failed to remove intermediate PDF {pdf_path}")

    return images, warnings


def shutil_which(cmd: str) -> str | None:
    """Thin shim so tests can monkey-patch `shutil.which`."""
    import shutil
    return shutil.which(cmd)


def _render_from_spec_code(spec_json: str, file_name: str) -> str:
    """Build a TypeScript program that renders *spec_json* via pipeshub-slides.

    The JSON is embedded as a literal in the program so the sandbox doesn't
    need to read a file. Each slide type maps to one of Deck's primitives;
    unknown types are skipped with a console warning so a single bad slide
    never breaks the whole render.
    """
    embedded_spec = json.dumps(spec_json)
    safe_file_name = json.dumps(file_name)
    return (
        "import { Deck, THEMES } from \"pipeshub-slides\";\n"
        "import * as path from \"path\";\n"
        "\n"
        f"const spec = JSON.parse({embedded_spec});\n"
        "const themeKey = spec.theme as keyof typeof THEMES;\n"
        "const theme = (THEMES as any)[themeKey] ?? THEMES.midnightExecutive;\n"
        "const deck = new Deck({\n"
        "  theme,\n"
        "  layout: (spec.layout as any) ?? \"16x9\",\n"
        "  title: spec.title,\n"
        "  author: spec.author,\n"
        "  company: spec.company,\n"
        "});\n"
        "\n"
        "for (const slide of spec.slides as any[]) {\n"
        "  switch (slide.type) {\n"
        "    case \"title\":\n"
        "      deck.titleSlide({\n"
        "        title: slide.title ?? \"\",\n"
        "        subtitle: slide.subtitle,\n"
        "        eyebrow: slide.eyebrow,\n"
        "        author: slide.author,\n"
        "        date: slide.date,\n"
        "      });\n"
        "      break;\n"
        "    case \"section_divider\":\n"
        "      deck.sectionDivider({\n"
        "        eyebrow: slide.eyebrow,\n"
        "        title: slide.title ?? \"\",\n"
        "        subtitle: slide.subtitle,\n"
        "      });\n"
        "      break;\n"
        "    case \"two_column\":\n"
        "      deck.twoColumn({\n"
        "        title: slide.title ?? \"\",\n"
        "        split: slide.split,\n"
        "        left: slide.left ?? { paragraph: \"\" },\n"
        "        right: slide.right ?? { paragraph: \"\" },\n"
        "      });\n"
        "      break;\n"
        "    case \"stat_grid\":\n"
        "      deck.statGrid({\n"
        "        title: slide.title ?? \"\",\n"
        "        stats: slide.stats ?? [],\n"
        "      });\n"
        "      break;\n"
        "    case \"icon_rows\":\n"
        "      deck.iconRows({\n"
        "        title: slide.title ?? \"\",\n"
        "        rows: slide.rows ?? [],\n"
        "      });\n"
        "      break;\n"
        "    case \"timeline\":\n"
        "      deck.timeline({\n"
        "        title: slide.title ?? \"\",\n"
        "        steps: slide.steps ?? [],\n"
        "      });\n"
        "      break;\n"
        "    case \"closing\":\n"
        "      deck.closing({\n"
        "        title: slide.title ?? \"\",\n"
        "        subtitle: slide.subtitle,\n"
        "        cta: slide.cta,\n"
        "      });\n"
        "      break;\n"
        "    default:\n"
        "      console.warn(`skipping unknown slide type: ${slide.type}`);\n"
        "  }\n"
        "}\n"
        "\n"
        f"const outPath = path.join(process.env.OUTPUT_DIR!, {safe_file_name});\n"
        "await deck.save(outPath);\n"
    )
